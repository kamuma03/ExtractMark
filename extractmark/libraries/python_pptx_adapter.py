"""LIB-07 -- python-pptx adapter."""

from __future__ import annotations

import logging
import time
from pathlib import Path

from extractmark.types import PageInput, PageOutput

logger = logging.getLogger(__name__)


class PythonPptxAdapter:
    """Adapter for python-pptx (LIB-07) -- native PPTX slide access."""

    lib_id = "LIB-07"

    def process_page(self, page: PageInput) -> PageOutput:
        return PageOutput(
            document_id=page.document_id,
            page_number=page.page_number,
            raw_text="",
            metadata={"note": "python-pptx works on PPTX files. Use process_document()."},
        )

    def process_document(self, doc_path: Path) -> list[PageOutput]:
        from pptx import Presentation
        from pptx.util import Inches

        document_id = doc_path.stem
        outputs: list[PageOutput] = []

        try:
            prs = Presentation(str(doc_path))
        except Exception as e:
            logger.error("python-pptx failed to open %s: %s", doc_path, e)
            return []

        for slide_num, slide in enumerate(prs.slides):
            start = time.perf_counter()
            text_parts: list[str] = []
            tables_md: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            text_parts.append(text)

                if shape.has_table:
                    table = shape.table
                    md_rows = []
                    for i, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        md_rows.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            md_rows.append("| " + " | ".join("---" for _ in cells) + " |")
                    if md_rows:
                        tables_md.append("\n".join(md_rows))

            elapsed_ms = (time.perf_counter() - start) * 1000
            full_text = "\n\n".join(text_parts)
            if tables_md:
                full_text += "\n\n" + "\n\n".join(tables_md)

            outputs.append(PageOutput(
                document_id=document_id,
                page_number=slide_num,
                raw_text=full_text,
                tables=tables_md,
                inference_time_ms=elapsed_ms,
                metadata={"library": self.lib_id, "slide": slide_num + 1},
            ))

        return outputs
